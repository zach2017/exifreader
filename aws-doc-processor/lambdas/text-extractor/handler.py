"""
Text Extractor Lambda
=====================
Triggered by: SQS (text-extract-queue)
              Step Functions (for PDF image extraction action)

Purpose:
  - Extracts text from documents: PDF (text layer), DOCX, RTF, HTML, and other formats
  - For Step Function 'extract_images' action: extracts embedded images from PDFs
  - Saves extracted text as .txt files to S3 extracted/ prefix
  - Updates DynamoDB metadata with extracted file references
  - Indexes extracted text content in Elasticsearch

Dependencies:
  - pdfplumber (PDF text extraction + image extraction)
  - python-docx (DOCX text extraction)
  - beautifulsoup4 (HTML text extraction)
  - openpyxl (Excel text extraction)
"""

import json
import os
import io
import logging
import tempfile
from datetime import datetime, timezone

import boto3

logger = logging.getLogger()
logger.setLevel(logging.INFO)

ENDPOINT_URL = os.environ.get('AWS_ENDPOINT_URL', None)
S3_BUCKET = os.environ.get('S3_BUCKET', 'docproc-bucket')
DYNAMODB_TABLE = os.environ.get('DYNAMODB_TABLE', 'document-metadata')
ES_URL = os.environ.get('ELASTICSEARCH_URL', 'http://localhost:9200')

boto_kwargs = {'endpoint_url': ENDPOINT_URL} if ENDPOINT_URL else {}
s3 = boto3.client('s3', **boto_kwargs)
dynamodb = boto3.resource('dynamodb', **boto_kwargs)
table = dynamodb.Table(DYNAMODB_TABLE)


def lambda_handler(event, context):
    """
    Main handler — dispatches based on event type.
    
    SQS Event: Extract text from the document referenced in the message.
    Step Function Event (action=extract_images): Extract images from a PDF.
    """
    logger.info(f"Text Extractor invoked: {json.dumps(event, default=str)[:500]}")
    
    # ── Step Function: Extract images from PDF ──
    if isinstance(event, dict) and event.get('action') == 'extract_images':
        return extract_images_from_pdf(event)
    
    # ── SQS trigger: Extract text from document ──
    if 'Records' in event:
        for record in event['Records']:
            body = json.loads(record.get('body', '{}'))
            process_text_extraction(body)
    elif 'file_id' in event:
        process_text_extraction(event)
    
    return {'statusCode': 200}


def process_text_extraction(message):
    """
    Extract text from a document based on its content type.
    
    Supported formats:
      - PDF: Extract text layer using pdfplumber
      - DOCX: Parse document.xml using python-docx
      - HTML: Strip tags using BeautifulSoup
      - XLSX: Read cell values using openpyxl
      - Other: Attempt raw text decode
    """
    file_id = message['file_id']
    s3_key = message['s3_key']
    content_type = message.get('content_type', '')
    
    logger.info(f"Extracting text: file_id={file_id}, type={content_type}")
    
    try:
        # Download file from S3 to temp directory
        with tempfile.NamedTemporaryFile(delete=False, suffix=_get_extension(s3_key)) as tmp:
            s3.download_file(S3_BUCKET, s3_key, tmp.name)
            tmp_path = tmp.name
        
        # Extract text based on content type
        if 'pdf' in content_type:
            text = extract_pdf_text(tmp_path)
        elif 'wordprocessingml' in content_type or content_type == 'application/msword':
            text = extract_docx_text(tmp_path)
        elif 'html' in content_type:
            text = extract_html_text(tmp_path)
        elif 'spreadsheetml' in content_type or 'ms-excel' in content_type:
            text = extract_xlsx_text(tmp_path)
        elif 'presentationml' in content_type or 'ms-powerpoint' in content_type:
            text = extract_pptx_text(tmp_path)
        else:
            text = extract_raw_text(tmp_path)
        
        if not text or not text.strip():
            logger.warning(f"No text extracted from {file_id}")
            text = "[No extractable text content found]"
        
        # Save extracted text to S3
        filename = s3_key.split('/')[-1]
        text_filename = os.path.splitext(filename)[0] + '.txt'
        extracted_key = f"extracted/{file_id}/{text_filename}"
        
        s3.put_object(
            Bucket=S3_BUCKET,
            Key=extracted_key,
            Body=text.encode('utf-8'),
            ContentType='text/plain'
        )
        logger.info(f"Saved extracted text: {extracted_key}")
        
        # Update DynamoDB
        now = datetime.now(timezone.utc).isoformat()
        table.update_item(
            Key={'file_id': file_id},
            UpdateExpression='''
                SET extracted_text_key = :tk, 
                    extracted_files = list_append(if_not_exists(extracted_files, :empty), :ef),
                    processing_steps = list_append(if_not_exists(processing_steps, :empty), :step)
            ''',
            ExpressionAttributeValues={
                ':tk': extracted_key,
                ':ef': [extracted_key],
                ':step': [f"{now}: Text extracted ({len(text)} chars) → {extracted_key}"],
                ':empty': []
            }
        )
        
        # Check if this was the last processing step — if so, mark COMPLETED
        check_and_complete(file_id)
        
        # Index in Elasticsearch
        index_in_elasticsearch(file_id, filename, text, s3_key, 'text_extraction')
        
        logger.info(f"Text extraction complete for {file_id}: {len(text)} chars")
        
    except Exception as e:
        logger.error(f"Text extraction failed for {file_id}: {e}", exc_info=True)
        update_status(file_id, 'ERROR')
        raise  # Re-raise for SQS retry
    
    finally:
        # Clean up temp file
        try:
            os.unlink(tmp_path)
        except Exception:
            pass


def extract_images_from_pdf(event):
    """
    Step Function action: Extract embedded images from a PDF.
    
    Downloads the PDF, extracts images page-by-page using pdfplumber,
    saves each image to S3, and returns the list of image S3 keys.
    
    Returns:
        dict with 'image_keys' list for the Step Function to route to OCR.
    """
    file_id = event['file_id']
    s3_key = event['s3_key']
    
    logger.info(f"Extracting images from PDF: {file_id}")
    
    image_keys = []
    
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as tmp:
            s3.download_file(S3_BUCKET, s3_key, tmp.name)
            tmp_path = tmp.name
        
        try:
            import pdfplumber
            
            with pdfplumber.open(tmp_path) as pdf:
                for page_num, page in enumerate(pdf.pages):
                    images = page.images
                    for img_idx, img in enumerate(images):
                        try:
                            # Extract image data from the page
                            img_name = f"page{page_num + 1}_img{img_idx + 1}.png"
                            img_key = f"extracted/{file_id}/images/{img_name}"
                            
                            # Get the image from the page
                            # pdfplumber provides image coordinates; we crop the page
                            x0 = img['x0']
                            top = img['top']
                            x1 = img['x1']
                            bottom = img['bottom']
                            
                            cropped = page.within_bbox((x0, top, x1, bottom))
                            if cropped:
                                img_obj = cropped.to_image(resolution=200)
                                
                                img_buffer = io.BytesIO()
                                img_obj.save(img_buffer, format='PNG')
                                img_buffer.seek(0)
                                
                                s3.put_object(
                                    Bucket=S3_BUCKET,
                                    Key=img_key,
                                    Body=img_buffer.read(),
                                    ContentType='image/png'
                                )
                                image_keys.append(img_key)
                                logger.info(f"Extracted image: {img_key}")
                                
                        except Exception as img_err:
                            logger.warning(f"Failed to extract image {img_idx} from page {page_num}: {img_err}")
                            continue
                            
        except ImportError:
            logger.warning("pdfplumber not available, attempting PyPDF2")
            # Fallback: use basic PDF library
            image_keys = extract_images_fallback(tmp_path, file_id)
        
        # Update DynamoDB with extracted images
        if image_keys:
            now = datetime.now(timezone.utc).isoformat()
            table.update_item(
                Key={'file_id': file_id},
                UpdateExpression='''
                    SET extracted_files = list_append(if_not_exists(extracted_files, :empty), :ef),
                        processing_steps = list_append(if_not_exists(processing_steps, :empty), :step)
                ''',
                ExpressionAttributeValues={
                    ':ef': image_keys,
                    ':step': [f"{now}: Extracted {len(image_keys)} images from PDF"],
                    ':empty': []
                }
            )
        
        os.unlink(tmp_path)
        
        return {
            'file_id': file_id,
            's3_key': s3_key,
            'image_keys': image_keys,
            'image_count': len(image_keys)
        }
        
    except Exception as e:
        logger.error(f"PDF image extraction failed: {e}", exc_info=True)
        return {
            'file_id': file_id,
            's3_key': s3_key,
            'image_keys': [],
            'image_count': 0,
            'error': str(e)
        }


def extract_images_fallback(pdf_path, file_id):
    """Fallback image extraction using PyPDF2 (fewer features but lighter)."""
    image_keys = []
    try:
        from PyPDF2 import PdfReader
        reader = PdfReader(pdf_path)
        
        img_count = 0
        for page_num, page in enumerate(reader.pages):
            if '/XObject' in page['/Resources']:
                xobjects = page['/Resources']['/XObject'].get_object()
                for obj_name in xobjects:
                    obj = xobjects[obj_name].get_object()
                    if obj['/Subtype'] == '/Image':
                        img_count += 1
                        img_key = f"extracted/{file_id}/images/page{page_num + 1}_img{img_count}.png"
                        
                        data = obj.get_data()
                        s3.put_object(
                            Bucket=S3_BUCKET,
                            Key=img_key,
                            Body=data,
                            ContentType='image/png'
                        )
                        image_keys.append(img_key)
    except Exception as e:
        logger.warning(f"Fallback image extraction failed: {e}")
    
    return image_keys


# ─── Format-Specific Text Extractors ─────────────────────────

def extract_pdf_text(filepath):
    """Extract text from PDF using pdfplumber (handles complex layouts)."""
    try:
        import pdfplumber
        text_parts = []
        with pdfplumber.open(filepath) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text_parts.append(page_text)
        return '\n\n'.join(text_parts)
    except ImportError:
        logger.warning("pdfplumber not available")
        return extract_raw_text(filepath)


def extract_docx_text(filepath):
    """Extract text from Word DOCX files using python-docx."""
    try:
        from docx import Document
        doc = Document(filepath)
        
        text_parts = []
        # Extract paragraphs
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)
        
        # Extract tables
        for tbl in doc.tables:
            for row in tbl.rows:
                row_text = ' | '.join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    text_parts.append(row_text)
        
        return '\n'.join(text_parts)
    except ImportError:
        logger.warning("python-docx not available")
        return extract_raw_text(filepath)


def extract_html_text(filepath):
    """Extract text from HTML files using BeautifulSoup."""
    try:
        from bs4 import BeautifulSoup
        with open(filepath, 'r', encoding='utf-8', errors='replace') as f:
            soup = BeautifulSoup(f.read(), 'html.parser')
            # Remove script and style elements
            for tag in soup(['script', 'style', 'nav', 'footer', 'header']):
                tag.decompose()
            return soup.get_text(separator='\n', strip=True)
    except ImportError:
        return extract_raw_text(filepath)


def extract_xlsx_text(filepath):
    """Extract text from Excel files using openpyxl."""
    try:
        from openpyxl import load_workbook
        wb = load_workbook(filepath, read_only=True, data_only=True)
        text_parts = []
        
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            text_parts.append(f"--- Sheet: {sheet_name} ---")
            for row in ws.iter_rows(values_only=True):
                row_text = ' | '.join(str(cell) for cell in row if cell is not None)
                if row_text:
                    text_parts.append(row_text)
        
        wb.close()
        return '\n'.join(text_parts)
    except ImportError:
        return extract_raw_text(filepath)


def extract_pptx_text(filepath):
    """Extract text from PowerPoint files."""
    try:
        from pptx import Presentation
        prs = Presentation(filepath)
        text_parts = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            text_parts.append(f"--- Slide {slide_num} ---")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            text_parts.append(para.text)
        
        return '\n'.join(text_parts)
    except ImportError:
        return extract_raw_text(filepath)


def extract_raw_text(filepath):
    """Last resort: attempt to read file as raw text."""
    try:
        with open(filepath, 'r', encoding='utf-8', errors='replace') as f:
            return f.read()
    except Exception:
        with open(filepath, 'rb') as f:
            raw = f.read()
            return raw.decode('utf-8', errors='replace')


# ─── Helper Functions ────────────────────────────────────────

def _get_extension(s3_key):
    """Get file extension from S3 key."""
    name = s3_key.split('/')[-1]
    if '.' in name:
        return '.' + name.rsplit('.', 1)[1]
    return ''


def check_and_complete(file_id):
    """
    Check if all processing is done for this file and mark as COMPLETED.
    For simplicity, we mark completed after text extraction.
    In production, you'd track expected vs completed tasks.
    """
    try:
        response = table.get_item(Key={'file_id': file_id})
        item = response.get('Item', {})
        
        # Simple heuristic: if we have extracted text, and status is still PROCESSING
        if item.get('status') == 'PROCESSING' and item.get('extracted_text_key'):
            # For non-PDF files, text extraction is the final step
            if item.get('file_category') != 'pdf':
                update_status(file_id, 'COMPLETED')
    except Exception as e:
        logger.warning(f"check_and_complete failed: {e}")


def update_status(file_id, status):
    """Update document status in DynamoDB."""
    try:
        table.update_item(
            Key={'file_id': file_id},
            UpdateExpression='SET #s = :s',
            ExpressionAttributeNames={'#s': 'status'},
            ExpressionAttributeValues={':s': status}
        )
    except Exception as e:
        logger.warning(f"Failed to update status: {e}")


def index_in_elasticsearch(file_id, filename, content, s3_key, source_type):
    """Index extracted text content in Elasticsearch."""
    try:
        import urllib.request
        doc = {
            'file_id': file_id,
            'filename': filename,
            'content': content[:50000],
            'file_type': source_type,
            'upload_time': datetime.now(timezone.utc).isoformat(),
            's3_key': s3_key,
            'source_type': source_type
        }
        req = urllib.request.Request(
            f"{ES_URL}/documents/_doc/{file_id}-text",
            data=json.dumps(doc).encode('utf-8'),
            headers={'Content-Type': 'application/json'},
            method='PUT'
        )
        urllib.request.urlopen(req, timeout=10)
        logger.info(f"Indexed text in Elasticsearch: {file_id}")
    except Exception as e:
        logger.warning(f"ES indexing failed (non-fatal): {e}")
